import os
from pathlib import Path
from typing import Optional
import logging
from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Inches
from PIL import Image

from app.utils.file_utils import ensure_directory_exists, get_file_size
from app.services.project_service import ProjectService

from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips, CompositeVideoClip

logger = logging.getLogger(__name__)


class ExportService:
    """PPT导出服务类"""
    
    def __init__(self, storage_dir: Optional[str] = None):
        """初始化导出服务
        
        Args:
            storage_dir: 存储目录，默认为当前工作目录下的storage/projects
        """
        if storage_dir is None:
            self.storage_dir = Path.cwd() / "storage" / "projects"
        else:
            self.storage_dir = Path(storage_dir)
        
        ensure_directory_exists(self.storage_dir)
    
    def set_16_9_slide_size(self, prs: Presentation):
        """设置PPT为16:9比例（标准宽高：13.333英寸 × 7.5英寸）"""
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    
    def add_centered_full_width_image(self, slide: Slide, width: Inches, height: Inches, img_path: Path):
        """在PPT正中央插入图片，宽度铺满PPT，高度按比例自适应"""
        # 获取PPT的宽高
        slide_width = width
        slide_height = height
        
        # 用PIL获取图片原始宽高比
        with Image.open(img_path) as img:
            img_width, img_height = img.size
            img_aspect_ratio = img_width / img_height
        
        # 计算图片高度（宽度=PPT宽度，保持原始宽高比）
        img_display_width = slide_width
        img_display_height = img_display_width / img_aspect_ratio
        
        # 计算居中位置（水平居中：left=0；垂直居中：top=(PPT高度-图片高度)/2）
        left = Inches(0)
        top = (slide_height - img_display_height) / 2
        
        # 插入图片
        slide.shapes.add_picture(
            str(img_path),
            left=left,
            top=top,
            width=img_display_width,
            height=img_display_height
        )
    
    def add_audio_to_top_left(self, slide: Slide, audio_path: Path):
        """在PPT左上角添加音频"""
        # 左上角位置（距离左边和上边各0.3英寸）
        left = Inches(0.3)
        top = Inches(0.3)
        
        width = Inches(1.5)
        height = Inches(1)
        
        # 添加音频对象
        try:
            audio = slide.shapes.add_movie(
                str(audio_path),
                left=left,
                top=top,
                width=width,
                height=height,
                mime_type="audio/mpeg"  # MP3格式用audio/mpeg，WAV用audio/wav
            )
            return audio
        except Exception as e:
            logger.warning(f"添加音频文件 {audio_path} 失败: {str(e)}")
            return None
    
    def get_sorted_image_files(self, project_id: str):
        """获取项目中按页码排序的图片文件列表"""
        images_dir = self.storage_dir / project_id / "images"
        if not images_dir.exists():
            return []
        
        # 获取所有PNG图片文件，使用3位数字格式
        image_files = list(images_dir.glob("page_*.png"))
        # 按页码排序，提取3位数字格式的页码
        image_files.sort(key=lambda x: int(x.stem.split('_')[1]))
        
        return image_files
    
    def get_sorted_audio_files(self, project_id: str):
        """获取项目中按页码排序的音频文件列表"""
        audio_dir = self.storage_dir / project_id / "audio"
        if not audio_dir.exists():
            return []
        
        # 获取所有MP3文件（页面合并音频）
        audio_files = list(audio_dir.glob("page_*.mp3"))
        # 按页码排序
        audio_files.sort(key=lambda x: int(x.stem.split('_')[1]))
        
        return audio_files
    
    def export_ppt(self, project_id: str, output_filename: Optional[str] = None) -> Path:
        """导出PPT
        
        Args:
            project_id: 项目ID
            output_filename: 输出文件名，如果为None则使用项目名
            
        Returns:
            导出的PPT文件路径
            
        Raises:
            ValueError: 当项目不存在或缺少必要文件时抛出异常
        """
        # 检查项目是否存在
        project_service = ProjectService()
        project = project_service.get_project(project_id)
        if not project:
            raise ValueError(f"项目不存在: {project_id}")
        
        # 检查项目是否有图片
        if not project.images:
            raise ValueError("项目未转换PDF为图片")
        
        # 获取图片和音频文件列表
        img_files = self.get_sorted_image_files(project_id)
        audio_files = self.get_sorted_audio_files(project_id)
        
        if not img_files:
            raise ValueError("项目图片目录为空")
        
        # 检查音频文件是否存在
        if not audio_files:
            logger.warning(f"项目 {project_id} 没有音频文件，将只导出图片")
        
        # 确保图片和音频文件数量一致（取较小值）
        num_slides = min(len(img_files), len(audio_files)) if audio_files else len(img_files)
        
        # 创建新的PPT
        prs = Presentation()
        
        # 设置PPT为16:9比例
        self.set_16_9_slide_size(prs)
        
        # 为每个图片和音频创建一个幻灯片
        for i in range(num_slides):
            # 添加空白幻灯片（布局6是空白布局）
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            
            # 插入居中且宽度铺满的图片
            img_path = img_files[i]
            self.add_centered_full_width_image(slide, prs.slide_width, prs.slide_height, img_path)
            
            # 如果有音频文件，插入左上角音频
            if audio_files and i < len(audio_files):
                audio_path = audio_files[i]
                self.add_audio_to_top_left(slide, audio_path)
            
            logger.info(f"已添加幻灯片 {i+1}: {img_path.name}")
        
        # 如果只有图片没有音频，也要添加剩余的图片幻灯片
        if len(img_files) > num_slides:
            for i in range(num_slides, len(img_files)):
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                
                img_path = img_files[i]
                self.add_centered_full_width_image(slide, prs.slide_width, prs.slide_height, img_path)
                
                logger.info(f"已添加幻灯片 {i+1}: {img_path.name}")
        
        # 确定输出文件名
        if output_filename is None:
            # 使用项目名作为文件名（去除特殊字符）
            safe_name = "".join(c for c in project.name if c.isalnum() or c in (' ', '-', '_')).rstrip()
            output_filename = f"{safe_name}.pptx"
        
        if not output_filename.endswith('.pptx'):
            output_filename += '.pptx'
        
        # 确保输出目录存在
        project_dir = self.storage_dir / project_id
        ensure_directory_exists(project_dir)
        
        # 保存PPT
        output_path = project_dir / output_filename
        prs.save(str(output_path))
        
        logger.info(f"PPT已导出完成！共生成 {len(img_files)} 张幻灯片，保存至 {output_path}")
        return output_path
    
    def get_ppt_file_path(self, project_id: str, filename: Optional[str] = None) -> Optional[Path]:
        """获取PPT文件路径
        
        Args:
            project_id: 项目ID
            filename: PPT文件名，如果为None则查找项目目录下的第一个pptx文件
            
        Returns:
            PPT文件路径，如果不存在则返回None
        """
        project_dir = self.storage_dir / project_id
        
        if filename:
            ppt_path = project_dir / filename
            return ppt_path if ppt_path.exists() else None
        else:
            # 查找项目目录下的第一个pptx文件
            ppt_files = list(project_dir.glob("*.pptx"))
            return ppt_files[0] if ppt_files else None
    
    def delete_ppt_file(self, project_id: str, filename: str) -> bool:
        """删除PPT文件
        
        Args:
            project_id: 项目ID
            filename: PPT文件名
            
        Returns:
            删除是否成功
        """
        try:
            ppt_path = self.get_ppt_file_path(project_id, filename)
            if ppt_path and ppt_path.exists():
                ppt_path.unlink()
                logger.info(f"已删除PPT文件: {ppt_path}")
                return True
            return False
        except Exception as e:
            logger.error(f"删除PPT文件失败: {str(e)}")
            return False

    def export_video(self, project_id: str, output_filename: Optional[str] = None, fps: int = 24) -> Path:
        """导出视频（将图片与音频合并为MP4视频）
        
        Args:
            project_id: 项目ID
            output_filename: 输出文件名，如果为None则使用项目名
            fps: 视频帧率，默认为24fps
            
        Returns:
            导出的视频文件路径
            
        Raises:
            ValueError: 当项目不存在或缺少必要文件时抛出异常
        """
        project_service = ProjectService()
        project = project_service.get_project(project_id)
        if not project:
            raise ValueError(f"项目不存在: {project_id}")
        
        if not project.images:
            raise ValueError("项目未转换PDF为图片")
        
        img_files = self.get_sorted_image_files(project_id)
        audio_files = self.get_sorted_audio_files(project_id)
        
        if not img_files:
            raise ValueError("项目图片目录为空")
        
        if not audio_files:
            logger.warning(f"项目 {project_id} 没有音频文件，将只导出图片为视频")
        
        video_clips = []
        
        if audio_files and len(audio_files) == len(img_files):
            for i, (img_path, audio_path) in enumerate(zip(img_files, audio_files)):
                try:
                    audio_clip = AudioFileClip(str(audio_path))
                    
                    img_clip = ImageClip(str(img_path))
                    
                    img_clip = img_clip.set_duration(audio_clip.duration)
                    img_clip = img_clip.set_audio(audio_clip)
                    
                    video_clips.append(img_clip)
                    logger.info(f"已处理第 {i+1} 页: {img_path.name} (音频时长: {audio_clip.duration:.2f}秒)")
                except Exception as e:
                    logger.warning(f"处理第 {i+1} 页失败: {str(e)}")
                    continue
        else:
            if audio_files:
                min_count = min(len(img_files), len(audio_files))
                logger.warning(f"图片数量({len(img_files)})与音频数量({len(audio_files)})不匹配，只处理前{min_count}页")
                
                for i in range(min_count):
                    try:
                        audio_path = audio_files[i]
                        audio_clip = AudioFileClip(str(audio_path))
                        
                        img_clip = ImageClip(str(img_files[i]))
                        img_clip = img_clip.set_duration(audio_clip.duration)
                        img_clip = img_clip.set_audio(audio_clip)
                        
                        video_clips.append(img_clip)
                        logger.info(f"已处理第 {i+1} 页")
                    except Exception as e:
                        logger.warning(f"处理第 {i+1} 页失败: {str(e)}")
                        continue
            else:
                logger.info("没有音频文件，为每张图片创建3秒的静态视频")
                for i, img_path in enumerate(img_files):
                    try:
                        img_clip = ImageClip(str(img_path)).set_duration(3)
                        video_clips.append(img_clip)
                        logger.info(f"已处理第 {i+1} 页: {img_path.name}")
                    except Exception as e:
                        logger.warning(f"处理第 {i+1} 页失败: {str(e)}")
                        continue
        
        if not video_clips:
            raise ValueError("没有可用的图片或音频来生成视频")
        
        if len(video_clips) == 1:
            final_video = video_clips[0]
        else:
            try:
                final_video = concatenate_videoclips(video_clips, method="compose")
            except Exception as e:
                logger.warning(f"使用compose方法拼接失败，尝试使用chain方法: {str(e)}")
                final_video = concatenate_videoclips(video_clips, method="chain")
        
        if output_filename is None:
            safe_name = "".join(c for c in project.name if c.isalnum() or c in (' ', '-', '_')).rstrip()
            output_filename = f"{safe_name}.mp4"
        
        if not output_filename.endswith('.mp4'):
            output_filename += '.mp4'
        
        project_dir = self.storage_dir / project_id
        ensure_directory_exists(project_dir)
        
        output_path = project_dir / output_filename
        
        try:
            final_video.write_videofile(
                str(output_path),
                fps=fps,
                codec='libx264',
                audio_codec='aac',
                preset='medium',
                threads=4,
                logger=None
            )
            logger.info(f"视频已导出完成！保存至 {output_path}")
        finally:
            final_video.close()
            for clip in video_clips:
                try:
                    clip.close()
                except Exception:
                    pass
        
        return output_path

    def get_video_file_path(self, project_id: str, filename: Optional[str] = None) -> Optional[Path]:
        """获取视频文件路径
        
        Args:
            project_id: 项目ID
            filename: 视频文件名，如果为None则查找项目目录下的第一个mp4文件
            
        Returns:
            视频文件路径，如果不存在则返回None
        """
        project_dir = self.storage_dir / project_id
        
        if filename:
            video_path = project_dir / filename
            return video_path if video_path.exists() else None
        else:
            video_files = list(project_dir.glob("*.mp4"))
            return video_files[0] if video_files else None